# loader.py
import os
import logging
from typing import List
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from concurrent.futures import ThreadPoolExecutor, as_completed

# ========================
# 模块初始化配置
# 配置日志输出
# ========================
# 将第三方库(httpcore, urllib3)的日志等级设为WARNING，避免输出过多调试信息
logging.getLogger("httpcore").setLevel(logging.WARNING)
logging.getLogger("urllib3").setLevel(logging.WARNING)

# 获取当前模块的日志记录器
logger = logging.getLogger(__name__)

# 基本日志配置，设置日志记录等级和输出格式
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)

"""
加载并读取TXT文本文件内容，返回文件中的全部文本。
如果读取过程中出现异常，会记录错误并返回空字符串。
"""


def load_txt(filepath):
    try:
        with open(filepath, "r", encoding="utf-8") as file:
            return file.read()
    except Exception as e:
        logging.error(f"读取TXT文件 {filepath} 失败: {e}")
        return ""


"""
使用PdfReader读取PDF文件中的文本内容。
将PDF的每个页面提取出的文本进行拼接并返回。
如果读取过程中出现异常，会记录错误并返回空字符串。
"""


def load_pdf(filepath):
    try:
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
        return text
    except Exception as e:
        logging.error(f"读取PDF文件 {filepath} 失败: {e}")
        return ""


"""
使用docx.Document读取DOCX文件中的段落文本。
将所有段落的文本用换行符拼接起来并返回。
如果读取过程中出现异常，会记录错误并返回空字符串。
"""


def load_docx(filepath):
    try:
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        logging.error(f"读取DOCX文件 {filepath} 失败: {e}")
        return ""


"""
使用pandas读取Excel文件，并将DataFrame转换为制表符分隔的字符串返回。
如果读取过程中出现异常，会记录错误并返回空字符串。
"""


def load_excel(filepath):
    try:
        df = pd.read_excel(filepath, engine='openpyxl')
        return df_to_text(df)
    except Exception as e:
        logging.error(f"读取Excel文件 {filepath} 失败: {e}")
        return ""


"""
使用pptx.Presentation读取PPTX文件。
遍历每一页(slide)和每个文本框(shape)的段落，将其文本内容按换行符拼接并返回。
如果读取过程中出现异常，会记录错误并返回空字符串。
"""


def load_pptx(filepath):
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + "\n"
        return text
    except Exception as e:
        logging.error(f"读取PPTX文件 {filepath} 失败: {e}")
        return ""


"""
将DataFrame转换为以制表符(\t)分隔的CSV格式字符串，并返回。
不包含索引(index)列。
"""


def df_to_text(df):
    return df.to_csv(index=False, sep='\t')


# 扩展名与对应的加载函数映射字典
# 根据文件后缀名选择合适的读取函数
LOADER_FUNCTIONS = {
    ".txt": load_txt,
    ".md": load_txt,  # Markdown 文件按文本文件处理，复用load_txt
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_excel,
    ".xls": load_excel,
    ".pptx": load_pptx,
}

"""
根据文件名的扩展名(ext)从LOADER_FUNCTIONS中查找对应的加载函数进行读取。
如果不支持的文件类型，则记录警告并返回空字符串。
"""


def load_file(filepath, filename):
    ext = os.path.splitext(filename)[1].lower()  # 提取文件后缀并转换为小写
    loader = LOADER_FUNCTIONS.get(ext)
    if loader:
        return loader(filepath)
    else:
        logging.warning(f"不支持的文件类型: {filename}")
        return ""


"""
从指定目录(directory)中批量加载所有文件的内容，并返回一个包含字典的列表。
每个字典结构为：{"filename": 文件名, "content": 文件内容}

其中使用线程池(ThreadPoolExecutor)并行地读取文件以提高效率。
"""


def load_documents(directory):
    docs = []
    files = os.listdir(directory)
    logging.info(f"在目录 {directory} 找到 {len(files)} 个文件")

    # 使用线程池并发加载文件，max_workers=4 表示最多4个并发线程
    with ThreadPoolExecutor(max_workers=4) as executor:
        future_to_filename = {}

        # 将每个文件的加载任务提交给线程池
        for filename in files:
            filepath = os.path.join(directory, filename)
            if os.path.isfile(filepath):  # 只处理文件，不处理子目录
                future = executor.submit(load_file, filepath, filename)
                future_to_filename[future] = filename

        # 收集加载结果
        for future in as_completed(future_to_filename):
            filename = future_to_filename[future]
            try:
                content = future.result()  # 获取文件读取结果
                if content:
                    # 只有当读取到非空内容时，才追加到docs列表
                    docs.append({"filename": filename, "content": content})
                else:
                    logging.info(f"文件 {filename} 没有加载到内容")
            except Exception as e:
                logging.error(f"加载文件 {filename} 时出错: {e}")

    return docs
